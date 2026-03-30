"""PPT 생성 모듈 — JSON 슬라이드 구조를 받아 python-pptx로 PPT를 만든다."""

import os
import re
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


# ── 레이아웃 상수 ───────────────────────────────────────────
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# 콘텐츠 슬라이드: 이미지 영역 (하단)
IMAGE_REGION_TOP = Inches(3.5)


def build_presentation(slide_data: dict, template_path: str, assets_dir: str) -> Presentation:
    """slide_data JSON을 받아 PPT를 생성한다."""
    prs = Presentation(template_path)

    # 템플릿의 기존 슬라이드 제거
    _remove_all_slides(prs)

    for slide_info in slide_data["slides"]:
        slide_type = slide_info.get("slide_type", "content")
        if slide_type == "title":
            _add_title_slide(prs, slide_info)
        elif slide_type == "closing":
            _add_closing_slide(prs, slide_info)
        else:
            _add_content_slide(prs, slide_info, assets_dir)

    return prs


# ── placeholder 접근 헬퍼 ─────────────────────────────────
def _get_placeholders(slide) -> dict:
    """slide.placeholders를 idx 기반 dict로 변환한다."""
    ph_dict = {}
    for ph in slide.placeholders:
        ph_dict[ph.placeholder_format.idx] = ph
    return ph_dict


# ── 슬라이드 제거 ──────────────────────────────────────────
def _remove_all_slides(prs: Presentation):
    """템플릿에 포함된 기존 슬라이드를 모두 제거한다."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


# ── 파일명 → 캡션 변환 ───────────────────────────────────
def _filename_to_caption(filename: str) -> str:
    """'figure1.png' → 'Figure 1', 'table2.png' → 'Table 2'"""
    name = os.path.splitext(filename)[0]
    match = re.match(r"(figure|fig|table)(\d+)", name, re.IGNORECASE)
    if match:
        label = match.group(1).capitalize()
        if label.lower().startswith("fig"):
            label = "Figure"
        num = match.group(2)
        return f"{label} {num}"
    return name


# ── 제목 슬라이드 ──────────────────────────────────────────
def _add_title_slide(prs: Presentation, info: dict):
    """레이아웃 0(제목 슬라이드)을 사용한 첫 페이지.
    Layout 0 placeholders:
      idx=13: Paper Title
      idx=14: Author, conference/journal name
      idx=15: Presenter name (우하단)
    """
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    ph_dict = _get_placeholders(slide)

    # 논문 제목
    if 13 in ph_dict:
        ph_dict[13].text = info.get("title", "")
        _style_placeholder_text(ph_dict[13], Pt(24), True)

    # 저자 + 학회
    if 14 in ph_dict:
        subtitle_parts = []
        if info.get("subtitle"):
            subtitle_parts.append(info["subtitle"])
        if info.get("venue"):
            subtitle_parts.append(info["venue"])
        if info.get("date"):
            subtitle_parts.append(info["date"])
        ph_dict[14].text = " | ".join(subtitle_parts)
        _style_placeholder_text(ph_dict[14], Pt(14), False)

    # 발표자 이름
    if 15 in ph_dict:
        ph_dict[15].text = info.get("presenter", "")
        _style_placeholder_text(ph_dict[15], Pt(14), False)

    # 하단에 생성 날짜 (흰색, 12pt)
    today_str = date.today().strftime("%Y-%m-%d")
    _add_text_box(slide, f"Generated: {today_str}",
                  Inches(0.5), Inches(6.8), Inches(3.0), Inches(0.4),
                  font_size=Pt(12), alignment=PP_ALIGN.LEFT,
                  color=RGBColor(0xFF, 0xFF, 0xFF))


# ── 콘텐츠 슬라이드 ───────────────────────────────────────
def _add_content_slide(prs: Presentation, info: dict, assets_dir: str):
    """레이아웃 1(제목 및 내용)을 사용한 본문 슬라이드."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    ph_dict = _get_placeholders(slide)

    # placeholder idx 13: 슬라이드 제목
    if 13 in ph_dict:
        ph_dict[13].text = info.get("title", "")
        _style_placeholder_text(ph_dict[13], Pt(24), True)

    # placeholder idx 14: 섹션 라벨
    if 14 in ph_dict:
        ph_dict[14].text = info.get("section", "")
        _style_placeholder_text(ph_dict[14], Pt(14), False,
                                color=RGBColor(0x44, 0x72, 0xC4))

    # placeholder idx 15: 불릿 내용
    if 15 in ph_dict:
        ph = ph_dict[15]
        bullets = info.get("bullets", [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                ph.text = bullet
                _style_paragraph(ph.text_frame.paragraphs[0], Pt(14))
            else:
                p = ph.text_frame.add_paragraph()
                p.text = bullet
                _style_paragraph(p, Pt(14))

    # 이미지 배치 (+ 캡션)
    asset_files = info.get("assets", [])
    asset_captions = info.get("asset_captions", {})
    if asset_files:
        _place_images(slide, asset_files, assets_dir, asset_captions)

    # 스피커 노트
    if info.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = info["speaker_notes"]


# ── 마무리 슬라이드 ───────────────────────────────────────
def _add_closing_slide(prs: Presentation, info: dict):
    """레이아웃 6(빈 화면)을 사용한 마무리 슬라이드."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    title = info.get("title", "Thank You!")
    subtitle = info.get("subtitle", "Q&A")

    _add_text_box(slide, title,
                  Inches(1.0), Inches(2.5), Inches(11.33), Inches(1.5),
                  font_size=Pt(40), bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, subtitle,
                  Inches(1.0), Inches(4.2), Inches(11.33), Inches(1.0),
                  font_size=Pt(24), alignment=PP_ALIGN.CENTER,
                  color=RGBColor(0x66, 0x66, 0x66))


# ── 이미지 배치 ───────────────────────────────────────────
def _place_images(slide, asset_files: list, assets_dir: str,
                   asset_captions: dict = None):
    """슬라이드 하단 영역에 이미지를 배치한다. 여러 장이면 가로로 나열.
    각 이미지 아래에 설명 캡션을 추가한다."""
    if asset_captions is None:
        asset_captions = {}

    valid = []  # (path, filename) 쌍
    for fname in asset_files:
        path = os.path.join(assets_dir, fname)
        if os.path.exists(path):
            valid.append((path, fname))

    if not valid:
        return

    n = len(valid)
    available_width = SLIDE_WIDTH_IN - 1.26  # 양쪽 마진
    per_image_width = (available_width - 0.3 * (n - 1)) / n
    max_height = 3.0  # 캡션 공간 확보를 위해 줄임

    for i, (img_path, fname) in enumerate(valid):
        left_in = 0.63 + i * (per_image_width + 0.3)
        img_actual_w, img_actual_h = _add_image(
            slide, img_path,
            Inches(left_in), IMAGE_REGION_TOP,
            Inches(per_image_width), Inches(max_height)
        )
        # 캡션: asset_captions에 있으면 사용, 없으면 파일명에서 생성
        caption = asset_captions.get(fname, _filename_to_caption(fname))
        caption_top = IMAGE_REGION_TOP + img_actual_h + Inches(0.05)
        _add_text_box(slide, caption,
                      Inches(left_in), caption_top,
                      Inches(per_image_width), Inches(0.4),
                      font_size=Pt(9), alignment=PP_ALIGN.CENTER,
                      color=RGBColor(0x66, 0x66, 0x66))


def _add_image(slide, image_path: str, left, top, max_width, max_height):
    """이미지를 비율 유지하며 최대 크기 내에 배치한다.
    실제 배치된 이미지의 (width, height)를 Emu 단위로 반환한다."""
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    aspect = img_w / img_h
    target_w = max_width
    target_h = int(target_w / aspect)

    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)

    # 영역 내 가운데 정렬
    offset_x = (max_width - target_w) // 2

    slide.shapes.add_picture(
        image_path,
        left + offset_x, top,
        target_w, target_h
    )

    return target_w, target_h


# ── 텍스트 유틸 ───────────────────────────────────────────
def _add_text_box(slide, text, left, top, width, height,
                  font_size=Pt(18), bold=False, alignment=PP_ALIGN.LEFT,
                  color=None):
    """텍스트 박스를 추가한다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _style_placeholder_text(placeholder, font_size, bold, color=None):
    """placeholder의 텍스트 스타일을 설정한다."""
    for paragraph in placeholder.text_frame.paragraphs:
        _style_paragraph(paragraph, font_size, bold, color)


def _style_paragraph(paragraph, font_size, bold=False, color=None):
    """paragraph의 모든 run에 스타일을 적용한다."""
    paragraph.space_before = Pt(4)
    paragraph.space_after = Pt(4)
    for run in paragraph.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
